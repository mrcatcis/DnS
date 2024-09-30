import os
from functools import cached_property
from io import BytesIO
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_PATTERN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.shapes.shapetree import Shape
from pptx.text.text import TextFrame
from pptx.slide import Slide
from pptx.util import Cm, Inches, Pt

from converter.spell import BaseSpell


class Card:
    WIDTH = Cm(6.365) 
    HEIGHT = Cm(8.89)
    BORDER_1_SIZE = Cm(0.3)
    
    TITLE_GAP_LEFT = Cm(0.25)
    TITLE_GAP_TOP = Cm(0.25)
    TITLE_GAP_BOTTOM = Cm(0.1)

    BODY_GAP_LEFT = Cm(0.4)
    BODY_GAP_TOP = Cm(0.9)
    

    def __init__(self, left: Cm, top: Cm, spell: BaseSpell):
        self.left = left
        self.top = top
        self.spell = spell
    

    @cached_property
    def font(self) -> str:
        return os.path.join(os.getcwd(), 'fonts/OpenSans-Regular.ttf')

    
    def add_to_slide(self, slide: Slide) -> None:
        self.add_border(slide)
        self.add_title(slide)
        self.add_body(slide)
        self.add_bottom_colontitul(slide)
    
    def add_border(self, slide: Slide) -> Shape:
        border: Shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.left,
            self.top,
            self.WIDTH,
            self.HEIGHT,
        )
        border.fill.patterned()
        border.fill.pattern = MSO_PATTERN.DOTTED_DIAMOND
        border.fill.back_color.rgb = RGBColor(230, 215, 200)
        border.fill.fore_color.rgb = RGBColor(150, 150, 150)
        border.line.color.rgb = RGBColor(0, 0, 0)
        border.line.width = self.BORDER_1_SIZE
        return border
    
    def add_title(self, slide: Slide) -> tuple[Shape, Shape]:
        LEVEL_SIZE = self.BODY_GAP_TOP - self.TITLE_GAP_BOTTOM

        title: Shape = slide.shapes.add_textbox(
            self.left + self.TITLE_GAP_LEFT,
            self.top + self.TITLE_GAP_TOP,
            self.WIDTH - 2.5 * self.TITLE_GAP_LEFT - LEVEL_SIZE,
            self.BODY_GAP_TOP - self.TITLE_GAP_BOTTOM - self.TITLE_GAP_TOP,
        )
        title.fill.background()
        title.line.color.rgb = RGBColor(0, 0, 0)
        title.line.width = Cm(0.05)
        title_tf = title.text_frame
        title_tf.text = f'{self.spell.name}'
        
        title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_tf.paragraphs[0].font.name = 'Open Sans'
        
        title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        level: Shape = slide.shapes.add_textbox(
            self.left + self.WIDTH - self.TITLE_GAP_LEFT - LEVEL_SIZE,
            # self.left + self.WIDTH - LEVEL_SPHERE,
            self.top + self.TITLE_GAP_TOP / 4,
            # self.top,
            LEVEL_SIZE,
            LEVEL_SIZE,
        )
        level.fill.background()
        # level.line.color.rgb = RGBColor(0, 0, 0)
        # level.line.width = Cm(0.05)
        level_tf = level.text_frame
        level_tf.text = f'{self.spell.level}'
        paragraph = level_tf.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.color.rgb = RGBColor(175, 0, 0)
        paragraph.font.name = 'Open Sans'
        paragraph.font.size = Pt(15)
        return title, level

    def add_body(self, slide: Slide) -> Shape:
        body: Shape = slide.shapes.add_textbox(
            self.left + self.BODY_GAP_LEFT,
            self.top + self.BODY_GAP_TOP,
            self.WIDTH - 2 * self.BODY_GAP_LEFT,
            self.HEIGHT - 2 * self.BODY_GAP_TOP,
        )
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(255, 255, 255)
        body.line.color.rgb = RGBColor(0, 0, 0)
        body.line.width = Cm(0.05)

        body_tf = body.text_frame
        body_tf.vertical_anchor = MSO_ANCHOR.TOP

        text = ''
        if self.spell.components.material:
            text += f'#!statsКомпоненты: {self.spell.components.material}\v'
        text += f'#!statsДлительность: {self.spell.duration}\v'
        text += f'#!statsАктивация: {self.spell.activation}\v'

        body_tf.text = text
        paragraph = body_tf.add_paragraph()
        for text_break in self.spell.description.split('\v'):
            for text_run in text_break.split('\n'):
                run = paragraph.add_run()
                run.text = text_run
            # if len(split_breaks) != 1:
            paragraph.add_line_break()

        body_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        body_tf.word_wrap = True
        body_tf._set_font('Open Sans', 10, False, False)


        self._enhance_body(body_tf)
        
        return body

    @staticmethod
    def _enhance_body(tf: TextFrame) -> None:
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                if run.text.startswith('#!bold'):
                    run.font.bold = True
                    run.text = run.text[len('#!bold'):]

                if run.text.startswith('#!stats'):
                    run.font.italic = True
                    run.text = run.text[len('#!stats'):]
                    run.font.size = Pt(6)
                
                if run.text.startswith('#!italic'):
                    run.font.italic = True
                    run.text = run.text[len('#!italic'):]

                if run.text.startswith('#!save_throw'):
                    run.font.color.rgb = RGBColor(254, 94, 0)
                    run.text = run.text[len('#!save_throw'):]
                
                if re.match(r'^#[\da-f]{6}', run.text):
                    run.font.color.rgb = RGBColor(int(run.text[1:3], 16), int(run.text[3:5], 16), int(run.text[5:7], 16))
                    run.text = run.text[len('#111111'):]



    def add_bottom_colontitul(self, slide: Slide) -> Shape:
        colontitul: Shape = slide.shapes.add_textbox(
            self.left + self.TITLE_GAP_LEFT,
            self.top + self.HEIGHT - self.BODY_GAP_TOP + self.TITLE_GAP_BOTTOM,
            self.WIDTH - 2 * self.TITLE_GAP_LEFT,
            self.BODY_GAP_TOP - self.TITLE_GAP_BOTTOM - self.TITLE_GAP_TOP
        )
        colontitul.fill.background()
        colontitul.line.color.rgb = RGBColor(0, 0, 0)
        colontitul.line.width = Cm(0.05)

        colontitul_tf = colontitul.text_frame
        components = ["В" if self.spell.components.vocal else "", "С" if self.spell.components.somatic else "", "М" if self.spell.components.material else ""]
        components = filter(bool, components)
        colontitul_tf.text = f'{self.spell.school.capitalize()} — {','.join(components)} — {self.spell.range} — {self.spell.source}'
        colontitul_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        colontitul_tf.word_wrap = True

        paragraph = colontitul_tf.paragraphs[0]
        paragraph.font.name = 'Open Sans'
        colontitul_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        return colontitul
    