from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import MSO_AUTO_SIZE

from converter.card import Card
from converter.spell import BaseSpell, TTGSpell, LssSpell
from converter.client import TTGCachedClient
from itertools import batched


class Converter:
    def __init__(self, spell_config_path: Path, spell_class: type[BaseSpell], cache_path: Path) -> None:
        self.ttg_client = TTGCachedClient(cache_path)

        self.spells = []
        self.spell_class = spell_class
        with open(spell_config_path, encoding='utf-8') as f:
            for batch in batched(filter(lambda x: not x.startswith('#'), f.readlines()), 9):
                spell_on_slide = []
                for spell in batch:
                    if spell.startswith('http'):
                        spell_on_slide.append(spell.rsplit('/', 1)[-1].strip())
                    else:
                        spell_on_slide.append(spell.strip())
                self.spells.append(spell_on_slide)            
 

    def create_presentation(self, path: Path):
        prs = Presentation()
        # A4
        prs.slide_height = Cm(29.7)
        prs.slide_width = Cm(21)
        blank_slide_layout = prs.slide_layouts[6]
        GAP = Cm(0.32)
        for spell_on_slide in self.spells:
            slide = prs.slides.add_slide(blank_slide_layout)

            left_position = (prs.slide_width - 3 * Card.WIDTH - 2 * GAP) / 2
            top_position = (prs.slide_height - 3 * Card.HEIGHT - 2 * GAP) / 2      
            for row in range(3):
                for col in range(3):
                    index = row * 3 + col
                    if index >= len(spell_on_slide):
                        break
                    spell_name = spell_on_slide[index]
                    spell = TTGSpell(self.ttg_client.get_spell(spell_name))
                    card = Card(
                        left_position + col * (Card.WIDTH + GAP),
                        top_position + row * (Card.HEIGHT + GAP),
                        spell,
                    )
                    card.add_to_slide(slide)

        prs.save(path)
